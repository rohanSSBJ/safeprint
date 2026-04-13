from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PPTX = "SafePrint_Pitch_Deck_Generated.pptx"
OUTPUT_MD = "pitch-deck-layout-spec.md"


THEME = {
    "fonts": {
        "title": "Aptos Display",
        "body": "Aptos",
        "fallback": "Calibri",
    },
    "colors": {
        "bg": RGBColor(255, 255, 255),
        "top_band": RGBColor(247, 243, 235),
        "title": RGBColor(36, 87, 255),
        "text": RGBColor(31, 41, 55),
        "muted": RGBColor(107, 114, 128),
        "line": RGBColor(209, 213, 219),
        "accent_fill": RGBColor(241, 245, 255),
        "accent_fill_2": RGBColor(236, 245, 255),
        "accent_strong": RGBColor(59, 130, 246),
        "success": RGBColor(40, 167, 69),
        "pilot": RGBColor(251, 146, 60),
        "live": RGBColor(139, 92, 246),
        "card_shadow": RGBColor(228, 232, 240),
    },
}


ASSETS = {
    "institute_logo": None,
    "startup_logo": None,
    "product_images": [None, None, None],
    "team_photos": [None, None, None, None, None],
    "faculty_photos": [None, None],
}


CONTENT = {
    "partner_tags": ["K-tech", "IIIT Dharwad", "Govt. of Karnataka", "Research Park", "Startup Karnataka"],
    "startup_name": "SafePrint",
    "tagline": "Privacy-first AI-powered campus print platform",
    "tech_tags": ["Privacy", "Campus Printing", "AI Operations"],
    "team_lead": "Team Lead: SafePrint Founding Team",
    "submission_date": "Pitching Submission Date: 13 Apr 2026",
    "problem": {
        "problem_statement": [
            "Students often print personal files on public shop systems where documents can be copied, cached, or forgotten.",
            "Campus-area print shops still operate manually, which creates queue uncertainty, payment friction, and weak trust.",
        ],
        "customer_segment": [
            "Primary users: college students printing IDs, certificates, forms, resumes, and financial documents.",
            "Secondary users: nearby print-shop owners who need a smoother order intake and better daily operations.",
        ],
        "market": [
            "TAM: student and academic document printing market across Indian campuses.",
            "SAM: tech-enabled campus and college-nearby print shops in target cities.",
            "SOM: first set of partner print shops and recurring student users acquired through pilot campuses.",
        ],
        "competition": [
            "Local print shops using walk-in or pen-drive based workflows.",
            "Ad-hoc WhatsApp, email, and Telegram-based print requests.",
            "Generic copy shops without privacy-first file handling or shop intelligence.",
        ],
    },
    "solution": {
        "solution_summary": [
            "SafePrint lets students discover nearby shops, upload documents securely, place print requests, and complete pickup with less waiting and more trust.",
            "For shop owners, SafePrint adds live queue visibility, demand insights, and a workflow that is easier to manage than manual file handling.",
        ],
        "benefits": [
            "Students save time, avoid repeated trips, and reduce the risk of sensitive files being left on public systems.",
            "Print shops handle more orders with better predictability, smoother intake, and fewer operational bottlenecks.",
        ],
        "differentiators": [
            "Privacy-first document handling is the core product promise rather than an afterthought.",
            "The platform combines secure upload, shop discovery, live status, payments, and owner-side intelligence in one flow.",
        ],
    },
    "assumptions": {
        "left": [
            "Students will prefer a trusted digital flow over sharing files through random chats or pen drives.",
            "Campus print shops are willing to adopt a lightweight dashboard if it improves throughput and trust.",
            "A browser-based workflow is enough for early pilots before deeper printer automation.",
        ],
        "right": [
            "We considered WhatsApp-only ordering, kiosk-first ordering, and fully manual shop-side coordination.",
            "The current approach won because it is easier to pilot, gives better privacy control, and creates reusable data for owner analytics.",
            "We can phase automation incrementally instead of forcing hardware-heavy adoption on day one.",
        ],
    },
    "bom": {
        "items": [
            ("Cloud infrastructure and storage", 160000),
            ("Development and integration overhead", 120000),
            ("Travel for shop onboarding and pilot coordination", 80000),
            ("Pilot operations, deployment, and support", 90000),
            ("Contingency and tooling", 50000),
        ],
        "total": 500000,
        "justification": [
            "Cloud spend covers application hosting, secure file handling, monitoring, and usage growth during pilots.",
            "Travel is necessary for onboarding print partners, validating operations on-site, and maintaining pilot momentum.",
            "Pilot allocation covers deployment support, troubleshooting, demos, and operational readiness before scale-up.",
        ],
    },
    "timeline": {
        "months": ["May", "Jun", "Jul", "Aug", "Sep", "Oct", "Nov", "Dec", "Jan", "Feb", "Mar"],
        "tasks": [
            ("Frontend experience", 0, 7, "dev"),
            ("Backend and shop APIs", 0, 7, "dev"),
            ("Security and print flow", 1, 7, "dev"),
            ("Dashboard and analytics", 2, 7, "dev"),
            ("Testing and hardening", 5, 7, "dev"),
            ("Pilot rollout", 8, 9, "pilot"),
            ("Go-live and early scale", 9, 10, "live"),
        ],
    },
    "team": {
        "members": [
            {"name": "Member 1", "subtitle": "Branch, Semester", "interest": "Product & UX"},
            {"name": "Member 2", "subtitle": "Branch, Semester", "interest": "Frontend"},
            {"name": "Member 3", "subtitle": "Branch, Semester", "interest": "Backend"},
            {"name": "Member 4", "subtitle": "Branch, Semester", "interest": "Security"},
            {"name": "Member 5", "subtitle": "Branch, Semester", "interest": "Analytics"},
        ],
        "faculty": [
            {"name": "Faculty Guide 1", "subtitle": "Branch", "interest": "Mentorship"},
            {"name": "Faculty Guide 2", "subtitle": "Branch", "interest": "Domain Guidance"},
        ],
    },
    "impact": [
        "SafePrint improves privacy for students by reducing dependence on public desktops and unmanaged file sharing.",
        "It helps local print shops become more digital, predictable, and easier to trust.",
        "Smarter routing and better queue visibility reduce wasted trips, idle waiting, and operating inefficiency.",
    ],
    "patent": [
        "Potential IP may emerge around privacy-preserving print request orchestration, secure document lifecycle handling, and shop-side release flows.",
        "The project also has scope for publication or conference presentation around campus micro-commerce and operational digitization.",
    ],
    "startup": {
        "why": [
            "SafePrint solves a real campus pain point that combines privacy, convenience, and business enablement.",
            "The model can expand from one college cluster to multiple campuses with a repeatable partner onboarding playbook.",
        ],
        "first_customers": [
            "Start with pilot print shops around one or two campuses and onboard students through hostels, class reps, and launch demos.",
            "Use early partner offers, referral loops, and visible trust messaging to win the first 10 recurring student customers.",
        ],
    },
    "product_sections": [
        "Architecture Diagram Placeholder",
        "Workflow / Operations Diagram Placeholder",
        "Screenshots and Dashboard Placeholder",
    ],
    "thank_you": "Thank You",
    "contact": "safeprint.team@iiitdwd.ac.in",
}


SLIDES = [
    {"number": 1, "title": "Cover", "layout": "cover", "regions": [
        "Top centered partner-logo row",
        "Centered institute logo placeholder",
        "Centered startup title, subtitle, and three tech pills",
        "Centered startup logo placeholder",
        "Presenter and submission date near bottom center",
    ]},
    {"number": 2, "title": "Problem", "layout": "quadrant_grid", "regions": [
        "Top-left: problem statement and why it matters",
        "Top-right: customer segment and example users",
        "Bottom-left: TAM, SAM, SOM framing",
        "Bottom-right: competition and existing alternatives",
    ]},
    {"number": 3, "title": "Solution and Differentiation", "layout": "two_column_split", "regions": [
        "Left: solution summary",
        "Right: economic benefit and differentiation",
    ]},
    {"number": 4, "title": "Product", "layout": "stacked_bands", "regions": [
        "Three centered visual bands for architecture, workflow, and screenshots",
    ]},
    {"number": 5, "title": "Assumptions and Alternatives", "layout": "two_column_split", "regions": [
        "Left: product and adoption assumptions",
        "Right: alternatives considered and rationale for chosen approach",
    ]},
    {"number": 6, "title": "BoM and Investment Ask", "layout": "two_column_split", "regions": [
        "Left: line-item BoM and investment ask summary",
        "Right: cost justification",
    ]},
    {"number": 7, "title": "Implementation Timeline", "layout": "gantt", "regions": [
        "Single large timeline region with month-by-month Gantt chart",
    ]},
    {"number": 8, "title": "Team", "layout": "team_grid", "regions": [
        "Top row of five profile cards",
        "Second row of two faculty guide cards aligned left",
    ]},
    {"number": 9, "title": "Social Economic Impact of the Product", "layout": "single_column_bullets", "regions": [
        "One main statement with supporting impact bullets",
    ]},
    {"number": 10, "title": "Patent and Publication Potential", "layout": "single_column_bullets", "regions": [
        "One main bullet and supporting sub-points",
    ]},
    {"number": 11, "title": "Startup Potential", "layout": "two_column_split", "regions": [
        "Left: why this should become a startup",
        "Right: first-10-customers plan",
    ]},
    {"number": 12, "title": "Thank You", "layout": "thank_you", "regions": [
        "Centered partner-logo row, institute logo placeholder, thank-you text, and optional contact line",
    ]},
]


@dataclass
class Area:
    left: float
    top: float
    width: float
    height: float


def rgb_to_hex(color: RGBColor) -> str:
    return f"#{color[0]:02X}{color[1]:02X}{color[2]:02X}"


def format_inr(amount: int) -> str:
    value = str(amount)
    if len(value) <= 3:
        return value
    last_three = value[-3:]
    rest = value[:-3]
    chunks = []
    while len(rest) > 2:
        chunks.append(rest[-2:])
        rest = rest[:-2]
    if rest:
        chunks.append(rest)
    return ",".join(reversed(chunks)) + "," + last_three


def set_slide_background(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = THEME["colors"]["bg"]


def add_rect(slide, area: Area, fill_color: RGBColor, line_color: RGBColor | None = None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(
        shape_type, Inches(area.left), Inches(area.top), Inches(area.width), Inches(area.height)
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    line = shape.line
    if line_color:
        line.color.rgb = line_color
        line.width = Pt(1)
    else:
        line.fill.background()
    return shape


def add_line(slide, x1: float, y1: float, x2: float, y2: float, color: RGBColor, width_pt: float = 1.3):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = color
    line.line.width = Pt(width_pt)
    return line


def add_textbox(
    slide,
    area: Area,
    text: str = "",
    *,
    font_size: float = 18,
    bold: bool = False,
    color: RGBColor | None = None,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    font_name: str | None = None,
    vertical_anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(Inches(area.left), Inches(area.top), Inches(area.width), Inches(area.height))
    text_frame = box.text_frame
    text_frame.word_wrap = True
    text_frame.margin_left = Pt(6)
    text_frame.margin_right = Pt(6)
    text_frame.margin_top = Pt(4)
    text_frame.margin_bottom = Pt(4)
    text_frame.vertical_anchor = vertical_anchor
    paragraph = text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = align
    run = paragraph.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or THEME["colors"]["text"]
    run.font.name = font_name or THEME["fonts"]["body"]
    return box


def add_bullets(
    slide,
    area: Area,
    items: Iterable[str],
    *,
    font_size: float = 19,
    color: RGBColor | None = None,
    level: int = 0,
    spacing_before: float = 4,
):
    box = slide.shapes.add_textbox(Inches(area.left), Inches(area.top), Inches(area.width), Inches(area.height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(6)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    for idx, item in enumerate(items):
        paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        paragraph.text = item
        paragraph.level = level
        paragraph.bullet = True
        paragraph.space_after = Pt(10)
        paragraph.space_before = Pt(spacing_before if idx else 0)
        run = paragraph.runs[0]
        run.font.size = Pt(font_size)
        run.font.name = THEME["fonts"]["body"]
        run.font.color.rgb = color or THEME["colors"]["text"]
    return box


def add_placeholder_card(slide, area: Area, label: str, subtitle: str = ""):
    add_rect(slide, area, THEME["colors"]["accent_fill"], THEME["colors"]["line"], radius=True)
    add_textbox(
        slide,
        Area(area.left + 0.1, area.top + 0.17, area.width - 0.2, 0.35),
        label,
        font_size=18,
        bold=True,
        color=THEME["colors"]["title"],
        align=PP_ALIGN.CENTER,
    )
    if subtitle:
        add_textbox(
            slide,
            Area(area.left + 0.12, area.top + 0.52, area.width - 0.24, area.height - 0.6),
            subtitle,
            font_size=11,
            color=THEME["colors"]["muted"],
            align=PP_ALIGN.CENTER,
        )


def add_partner_row(slide, top: float, centered: bool = False) -> None:
    tags = CONTENT["partner_tags"]
    pill_width = 1.17
    gap = 0.12
    total_width = len(tags) * pill_width + (len(tags) - 1) * gap
    start_left = (13.333 - total_width) / 2 if centered else 7.15
    for idx, tag in enumerate(tags):
        left = start_left + idx * (pill_width + gap)
        pill = add_rect(
            slide,
            Area(left, top, pill_width, 0.36),
            RGBColor(255, 255, 255),
            THEME["colors"]["line"],
            radius=True,
        )
        pill.shadow.inherit = False
        add_textbox(
            slide,
            Area(left + 0.03, top + 0.035, pill_width - 0.06, 0.28),
            tag,
            font_size=11,
            bold=True,
            color=THEME["colors"]["muted"],
            align=PP_ALIGN.CENTER,
        )


def add_common_shell(slide, title: str, number: int) -> None:
    set_slide_background(slide)
    add_rect(slide, Area(0, 0, 13.333, 0.54), THEME["colors"]["top_band"])
    add_textbox(
        slide,
        Area(0.18, 0.1, 2.9, 0.32),
        "Institute Logo",
        font_size=18,
        color=THEME["colors"]["title"],
    )
    add_partner_row(slide, top=0.08, centered=False)
    add_textbox(
        slide,
        Area(0.22, 0.86, 5.0, 0.55),
        title,
        font_size=28,
        bold=False,
        color=THEME["colors"]["title"],
        font_name=THEME["fonts"]["title"],
    )
    add_line(slide, 0.52, 7.14, 12.86, 7.14, THEME["colors"]["line"], 1.6)
    add_textbox(
        slide,
        Area(2.1, 7.18, 8.9, 0.22),
        "© Copyright - IIIT Dharwad Research Park Foundation - Not to be reproduced without permission from author",
        font_size=10.5,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Area(12.97, 7.15, 0.2, 0.22),
        str(number),
        font_size=10,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.RIGHT,
    )


def add_logo_placeholder(slide, area: Area, label: str) -> None:
    add_rect(slide, area, THEME["colors"]["accent_fill_2"], THEME["colors"]["line"], radius=True)
    add_textbox(
        slide,
        Area(area.left + 0.08, area.top + 0.09, area.width - 0.16, area.height - 0.18),
        label,
        font_size=18,
        color=THEME["colors"]["title"],
        align=PP_ALIGN.CENTER,
        vertical_anchor=MSO_ANCHOR.MIDDLE,
    )


def add_picture_or_placeholder(slide, area: Area, asset_path: str | None, label: str) -> None:
    if asset_path and Path(asset_path).exists():
        slide.shapes.add_picture(asset_path, Inches(area.left), Inches(area.top), Inches(area.width), Inches(area.height))
    else:
        add_placeholder_card(slide, area, label, "Replace this placeholder with the final asset.")


def build_cover(slide) -> None:
    set_slide_background(slide)
    add_partner_row(slide, top=0.86, centered=True)
    add_logo_placeholder(slide, Area(5.35, 1.85, 2.65, 0.72), "Institute Logo")
    add_textbox(
        slide,
        Area(3.1, 3.15, 7.2, 0.65),
        CONTENT["startup_name"],
        font_size=36,
        bold=True,
        color=THEME["colors"]["text"],
        align=PP_ALIGN.CENTER,
        font_name=THEME["fonts"]["title"],
    )
    add_textbox(
        slide,
        Area(2.2, 3.72, 8.95, 0.4),
        CONTENT["tagline"],
        font_size=20,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    pill_top = 4.18
    pill_width = 1.65
    gap = 0.18
    total = len(CONTENT["tech_tags"]) * pill_width + (len(CONTENT["tech_tags"]) - 1) * gap
    start = (13.333 - total) / 2
    for idx, pill_text in enumerate(CONTENT["tech_tags"]):
        left = start + idx * (pill_width + gap)
        add_rect(slide, Area(left, pill_top, pill_width, 0.38), THEME["colors"]["text"], None, radius=True)
        add_textbox(
            slide,
            Area(left + 0.02, pill_top + 0.03, pill_width - 0.04, 0.26),
            pill_text,
            font_size=13,
            color=RGBColor(255, 255, 255),
            align=PP_ALIGN.CENTER,
        )
    add_logo_placeholder(slide, Area(5.48, 5.22, 2.36, 0.72), "Startup Logo")
    add_textbox(
        slide,
        Area(4.25, 6.02, 4.9, 0.28),
        CONTENT["team_lead"],
        font_size=17,
        color=THEME["colors"]["text"],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Area(4.0, 6.34, 5.35, 0.28),
        CONTENT["submission_date"],
        font_size=16,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, 0.52, 7.14, 12.86, 7.14, THEME["colors"]["line"], 1.6)
    add_textbox(
        slide,
        Area(2.1, 7.18, 8.9, 0.22),
        "© Copyright - IIIT Dharwad Research Park Foundation - Not to be reproduced without permission from author",
        font_size=10.5,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Area(12.97, 7.15, 0.2, 0.22),
        "1",
        font_size=10,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.RIGHT,
    )


def build_problem(slide) -> None:
    add_common_shell(slide, "Problem", 2)
    add_line(slide, 6.67, 2.2, 6.67, 6.35, THEME["colors"]["muted"], 1.2)
    add_line(slide, 1.65, 4.33, 12.02, 4.33, THEME["colors"]["muted"], 1.2)
    add_bullets(slide, Area(2.0, 2.35, 3.95, 1.7), CONTENT["problem"]["problem_statement"], font_size=18)
    add_bullets(slide, Area(7.15, 2.52, 4.2, 1.6), CONTENT["problem"]["customer_segment"], font_size=18)
    add_bullets(slide, Area(2.0, 4.98, 3.95, 1.2), CONTENT["problem"]["market"], font_size=18)
    add_bullets(slide, Area(7.15, 4.78, 4.25, 1.45), CONTENT["problem"]["competition"], font_size=18)


def build_two_column(slide, title: str, number: int, left_items: list[str], right_items: list[str], left_width=3.55):
    add_common_shell(slide, title, number)
    divider_x = 6.57
    add_line(slide, divider_x, 2.45, divider_x, 5.85, THEME["colors"]["muted"], 1.2)
    add_bullets(slide, Area(2.75, 3.05, left_width, 2.15), left_items, font_size=20)
    add_bullets(slide, Area(6.95, 2.42, 4.45, 2.75), right_items, font_size=19)


def build_product(slide) -> None:
    add_common_shell(slide, "Product", 4)
    sections = CONTENT["product_sections"]
    tops = [2.08, 3.2, 4.35]
    for idx, label in enumerate(sections):
        add_picture_or_placeholder(slide, Area(3.0, tops[idx], 7.2, 0.82), ASSETS["product_images"][idx], label)


def build_bom(slide) -> None:
    add_common_shell(slide, "BoM and Investment Ask", 6)
    add_line(slide, 6.75, 2.42, 6.75, 5.88, THEME["colors"]["muted"], 1.2)
    left_box = add_rect(
        slide, Area(0.82, 2.65, 5.45, 2.8), RGBColor(255, 255, 255), THEME["colors"]["line"], radius=True
    )
    left_box.shadow.inherit = False
    bullet_items = [f"{name}: INR {format_inr(amount)}" for name, amount in CONTENT["bom"]["items"]]
    add_bullets(slide, Area(0.95, 2.8, 5.0, 1.95), bullet_items, font_size=16)
    total_card = add_rect(
        slide, Area(1.05, 4.88, 4.35, 0.72), THEME["colors"]["accent_fill"], THEME["colors"]["accent_strong"], radius=True
    )
    total_card.line.width = Pt(1.8)
    add_textbox(
        slide,
        Area(1.18, 5.0, 4.1, 0.35),
        f"Investment Ask: INR {format_inr(CONTENT['bom']['total'])}",
        font_size=20,
        bold=True,
        color=THEME["colors"]["title"],
        align=PP_ALIGN.CENTER,
    )
    add_bullets(slide, Area(7.4, 2.68, 4.25, 2.45), CONTENT["bom"]["justification"], font_size=18)


def build_gantt(slide) -> None:
    add_common_shell(slide, "Implementation Timeline", 7)
    canvas = Area(1.0, 2.0, 11.2, 3.85)
    add_rect(slide, canvas, RGBColor(252, 253, 255), THEME["colors"]["line"], radius=True)
    label_x = 2.6
    months = CONTENT["timeline"]["months"]
    tasks = CONTENT["timeline"]["tasks"]
    month_start = 3.45
    month_width = 0.72
    top = 2.38
    add_textbox(slide, Area(1.28, top, 1.8, 0.28), "Workstream", font_size=13, bold=True, color=THEME["colors"]["muted"])
    for idx, month in enumerate(months):
        left = month_start + idx * month_width
        add_rect(slide, Area(left, top - 0.02, month_width, 0.34), THEME["colors"]["accent_fill"], THEME["colors"]["line"])
        add_textbox(
            slide,
            Area(left, top + 0.02, month_width, 0.22),
            month,
            font_size=11,
            bold=True,
            color=THEME["colors"]["muted"],
            align=PP_ALIGN.CENTER,
        )
    row_top = 2.8
    row_height = 0.42
    colors = {
        "dev": THEME["colors"]["accent_strong"],
        "pilot": THEME["colors"]["pilot"],
        "live": THEME["colors"]["live"],
    }
    legend = [("Development", "dev"), ("Pilot", "pilot"), ("Go-live", "live")]
    for idx, (label, key) in enumerate(legend):
        left = 8.65 + idx * 1.1
        add_rect(slide, Area(left, 5.95, 0.16, 0.16), colors[key], None, radius=True)
        add_textbox(slide, Area(left + 0.2, 5.89, 0.8, 0.22), label, font_size=10.5, color=THEME["colors"]["muted"])
    for row_index, (label, start, end, kind) in enumerate(tasks):
        current_top = row_top + row_index * row_height
        add_line(slide, 1.25, current_top + 0.27, 12.0, current_top + 0.27, THEME["colors"]["line"], 0.7)
        add_textbox(
            slide,
            Area(1.28, current_top + 0.03, label_x - 1.5, 0.22),
            label,
            font_size=13,
            color=THEME["colors"]["text"],
        )
        bar_left = month_start + start * month_width + 0.04
        bar_width = (end - start + 1) * month_width - 0.08
        bar = add_rect(slide, Area(bar_left, current_top + 0.06, bar_width, 0.2), colors[kind], None, radius=True)
        bar.shadow.inherit = False
    add_textbox(
        slide,
        Area(1.3, 6.02, 5.8, 0.24),
        "May-Dec: development | Jan-Feb: pilot rollout | Feb-Mar: go-live and early scale",
        font_size=11,
        color=THEME["colors"]["muted"],
    )


def add_profile_card(slide, area: Area, name: str, subtitle: str, interest: str) -> None:
    card = add_rect(slide, area, RGBColor(255, 255, 255), THEME["colors"]["line"], radius=True)
    card.shadow.inherit = False
    circle = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        Inches(area.left + area.width / 2 - 0.45),
        Inches(area.top + 0.12),
        Inches(0.9),
        Inches(0.9),
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = THEME["colors"]["accent_fill"]
    circle.line.color.rgb = THEME["colors"]["line"]
    add_textbox(
        slide,
        Area(area.left + 0.15, area.top + 0.42, area.width - 0.3, 0.16),
        "Photo",
        font_size=11,
        bold=True,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Area(area.left + 0.1, area.top + 1.05, area.width - 0.2, 0.22),
        name,
        font_size=14,
        bold=True,
        color=THEME["colors"]["text"],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Area(area.left + 0.1, area.top + 1.26, area.width - 0.2, 0.18),
        subtitle,
        font_size=12,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Area(area.left + 0.1, area.top + 1.46, area.width - 0.2, 0.24),
        interest,
        font_size=11.5,
        color=THEME["colors"]["text"],
        align=PP_ALIGN.CENTER,
    )


def build_team(slide) -> None:
    add_common_shell(slide, "Team", 8)
    start_left = 0.76
    top_row_top = 2.12
    card_w = 2.28
    card_h = 1.96
    gap = 0.19
    for idx, member in enumerate(CONTENT["team"]["members"]):
        left = start_left + idx * (card_w + gap)
        add_profile_card(slide, Area(left, top_row_top, card_w, card_h), member["name"], member["subtitle"], member["interest"])
    bottom_top = 4.72
    for idx, guide in enumerate(CONTENT["team"]["faculty"]):
        left = start_left + idx * (card_w + gap)
        add_profile_card(slide, Area(left, bottom_top, card_w, card_h), guide["name"], guide["subtitle"], guide["interest"])


def build_single_column(slide, title: str, number: int, items: list[str]) -> None:
    add_common_shell(slide, title, number)
    add_bullets(slide, Area(0.9, 3.0, 10.7, 1.8), items, font_size=20)


def build_thank_you(slide) -> None:
    set_slide_background(slide)
    add_partner_row(slide, top=0.86, centered=True)
    add_logo_placeholder(slide, Area(5.35, 1.85, 2.65, 0.72), "Institute Logo")
    add_textbox(
        slide,
        Area(4.15, 3.84, 5.0, 0.7),
        CONTENT["thank_you"],
        font_size=34,
        bold=True,
        color=THEME["colors"]["text"],
        align=PP_ALIGN.CENTER,
        font_name=THEME["fonts"]["title"],
    )
    add_textbox(
        slide,
        Area(4.12, 4.56, 5.1, 0.28),
        CONTENT["contact"],
        font_size=15,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_line(slide, 0.52, 7.14, 12.86, 7.14, THEME["colors"]["line"], 1.6)
    add_textbox(
        slide,
        Area(2.1, 7.18, 8.9, 0.22),
        "© Copyright - IIIT Dharwad Research Park Foundation - Not to be reproduced without permission from author",
        font_size=10.5,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.CENTER,
    )
    add_textbox(
        slide,
        Area(12.97, 7.15, 0.2, 0.22),
        "12",
        font_size=10,
        color=THEME["colors"]["muted"],
        align=PP_ALIGN.RIGHT,
    )


def build_slide(prs: Presentation, slide_def: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    layout = slide_def["layout"]
    if layout == "cover":
        build_cover(slide)
    elif layout == "quadrant_grid":
        build_problem(slide)
    elif layout == "two_column_split" and slide_def["number"] == 3:
        right = CONTENT["solution"]["benefits"] + CONTENT["solution"]["differentiators"]
        build_two_column(slide, slide_def["title"], slide_def["number"], CONTENT["solution"]["solution_summary"], right)
    elif layout == "stacked_bands":
        build_product(slide)
    elif layout == "two_column_split" and slide_def["number"] == 5:
        build_two_column(slide, slide_def["title"], slide_def["number"], CONTENT["assumptions"]["left"], CONTENT["assumptions"]["right"])
    elif layout == "two_column_split" and slide_def["number"] == 6:
        build_bom(slide)
    elif layout == "gantt":
        build_gantt(slide)
    elif layout == "team_grid":
        build_team(slide)
    elif layout == "single_column_bullets" and slide_def["number"] == 9:
        build_single_column(slide, slide_def["title"], slide_def["number"], CONTENT["impact"])
    elif layout == "single_column_bullets" and slide_def["number"] == 10:
        build_single_column(slide, slide_def["title"], slide_def["number"], CONTENT["patent"])
    elif layout == "two_column_split" and slide_def["number"] == 11:
        build_two_column(slide, slide_def["title"], slide_def["number"], CONTENT["startup"]["why"], CONTENT["startup"]["first_customers"])
    elif layout == "thank_you":
        build_thank_you(slide)
    else:
        raise ValueError(f"Unhandled slide layout: {layout}")


def markdown_slide_content(slide_def: dict) -> list[str]:
    number = slide_def["number"]
    if number == 1:
        return [
            f"Startup title: {CONTENT['startup_name']}",
            f"Subtitle: {CONTENT['tagline']}",
            f"Tech pills: {', '.join(CONTENT['tech_tags'])}",
            f"Presenter line: {CONTENT['team_lead']}",
        ]
    if number == 2:
        return [
            f"Top-left: {' | '.join(CONTENT['problem']['problem_statement'])}",
            f"Top-right: {' | '.join(CONTENT['problem']['customer_segment'])}",
            f"Bottom-left: {' | '.join(CONTENT['problem']['market'])}",
            f"Bottom-right: {' | '.join(CONTENT['problem']['competition'])}",
        ]
    if number == 3:
        return [
            f"Left: {' | '.join(CONTENT['solution']['solution_summary'])}",
            f"Right: {' | '.join(CONTENT['solution']['benefits'] + CONTENT['solution']['differentiators'])}",
        ]
    if number == 4:
        return [f"Visual bands: {' | '.join(CONTENT['product_sections'])}"]
    if number == 5:
        return [
            f"Left: {' | '.join(CONTENT['assumptions']['left'])}",
            f"Right: {' | '.join(CONTENT['assumptions']['right'])}",
        ]
    if number == 6:
        items = "; ".join(f"{name} - INR {format_inr(amount)}" for name, amount in CONTENT["bom"]["items"])
        return [
            f"BoM items: {items}",
            f"Total investment ask: INR {format_inr(CONTENT['bom']['total'])}",
            f"Justification: {' | '.join(CONTENT['bom']['justification'])}",
        ]
    if number == 7:
        items = "; ".join(f"{name}: {CONTENT['timeline']['months'][start]} to {CONTENT['timeline']['months'][end]}" for name, start, end, _ in CONTENT["timeline"]["tasks"])
        return [f"Gantt tracks: {items}"]
    if number == 8:
        members = ", ".join(member["name"] for member in CONTENT["team"]["members"])
        faculty = ", ".join(guide["name"] for guide in CONTENT["team"]["faculty"])
        return [f"Top row members: {members}", f"Bottom row faculty: {faculty}"]
    if number == 9:
        return [f"Impact bullets: {' | '.join(CONTENT['impact'])}"]
    if number == 10:
        return [f"Patent/publication bullets: {' | '.join(CONTENT['patent'])}"]
    if number == 11:
        return [
            f"Left: {' | '.join(CONTENT['startup']['why'])}",
            f"Right: {' | '.join(CONTENT['startup']['first_customers'])}",
        ]
    return [f"Message: {CONTENT['thank_you']}", f"Contact: {CONTENT['contact']}"]


def build_markdown() -> str:
    lines = [
        "# SafePrint Pitch Deck Layout Spec",
        "",
        "This file is generated by `generate_pitch_deck.py` and documents the produced PowerPoint deck.",
        "The deck follows the IIIT Dharwad NAIN template geometry while using a modern light-theme visual system.",
        "",
        "## Theme",
        "",
        "- Slide size: 16:9 widescreen approximating the original 1920x1080 template",
        f"- Title color: `{rgb_to_hex(THEME['colors']['title'])}`",
        f"- Body text color: `{rgb_to_hex(THEME['colors']['text'])}`",
        f"- Divider color: `{rgb_to_hex(THEME['colors']['line'])}`",
        "- Shared styling: pale warm top band, clean footer, rounded cards, light placeholders",
        "",
    ]
    for slide_def in SLIDES:
        lines.append(f"## Slide {slide_def['number']} - {slide_def['title']}")
        lines.append("")
        lines.append(f"- Layout type: `{slide_def['layout']}`")
        lines.append("- Region breakdown:")
        for region in slide_def["regions"]:
            lines.append(f"  - {region}")
        lines.append("- SafePrint content inserted:")
        for item in markdown_slide_content(slide_def):
            lines.append(f"  - {item}")
        lines.append("- Placeholder notes:")
        if slide_def["number"] in (1, 12):
            lines.append("  - Institute and startup logo regions use polished placeholders when no asset path is configured.")
        elif slide_def["number"] == 4:
            lines.append("  - Product bands render placeholder cards unless product diagrams or screenshots are provided.")
        elif slide_def["number"] == 8:
            lines.append("  - Team and faculty cards render circular photo placeholders when images are missing.")
        else:
            lines.append("  - This slide is fully generated from text and shape primitives, so no external assets are required.")
        lines.append("")
    return "\n".join(lines)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide_def in SLIDES:
        build_slide(prs, slide_def)
    return prs


def main() -> None:
    prs = build_presentation()
    prs.save(OUTPUT_PPTX)
    Path(OUTPUT_MD).write_text(build_markdown(), encoding="utf-8")
    print(f"Generated {OUTPUT_PPTX}")
    print(f"Generated {OUTPUT_MD}")


if __name__ == "__main__":
    main()
